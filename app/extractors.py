"""
Local document extractors for lang_un_rag.

Each extractor returns a list of dicts: {"text": str, "metadata": {...}}

Supports:
- PDF (born-digital via pdfplumber/pypdf; fallback to OCR if very little text)
- DOCX
- PPTX
- HTML
- Images (OCR via Tesseract)
- CSV
- Plain text fallback
"""
import os
import mimetypes
import warnings
import logging
from typing import List, Dict

# Setup logger for extractors
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# Suppress specific PDF parsing warnings
warnings.filterwarnings("ignore", message=".*FontBBox.*")
warnings.filterwarnings("ignore", message=".*cannot be parsed as 4 floats.*")

# Configure logging to reduce PDF parsing noise
logging.getLogger("pdfplumber").setLevel(logging.ERROR)
logging.getLogger("pypdf").setLevel(logging.ERROR)

try:
    import magic  # python-magic
except Exception:
    magic = None

# PDF
import pdfplumber
from pypdf import PdfReader

# DOCX
import docx

# PPTX
from pptx import Presentation

# HTML
from bs4 import BeautifulSoup

# Images / OCR
from PIL import Image
import pytesseract
from pdf2image import convert_from_path

# Email
import email
from email import policy
from email.parser import BytesParser
import re

# CSV/Excel
import pandas as pd

# Audio transcription
import requests
import nltk
from nltk.tokenize import sent_tokenize

# Download NLTK data on first import (punkt tokenizer for sentence splitting)
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)

# Heuristic thresholds
MIN_PDF_TEXT_LEN = 200  # if extracted text shorter than this, consider OCR fallback

def detect_mime(path: str) -> str:
    if magic:
        try:
            return magic.from_file(path, mime=True)
        except Exception:
            pass
    return mimetypes.guess_type(path)[0] or "application/octet-stream"

def extract_from_pdf(path: str) -> List[Dict]:
    texts = []
    
    # First try with pdfplumber with suppressed warnings
    try:
        # Suppress stderr temporarily to hide FontBBox warnings
        import sys
        import os
        from contextlib import redirect_stderr
        
        with open(os.devnull, 'w') as devnull:
            with redirect_stderr(devnull):
                with pdfplumber.open(path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        try:
                            text = page.extract_text() or ""
                            text = text.strip()
                            if text:
                                texts.append({"text": text, "metadata": {"file": os.path.basename(path), "page": i + 1, "file_type": "pdf"}})
                        except Exception as page_error:
                            # Skip problematic pages but continue with others
                            print(f"Warning: Failed to extract text from page {i + 1} of {os.path.basename(path)}: {page_error}")
                            continue
    except Exception as pdf_error:
        print(f"Warning: pdfplumber failed for {os.path.basename(path)}: {pdf_error}")
        # fallback to pypdf
        try:
            import sys
            import os
            from contextlib import redirect_stderr
            
            with open(os.devnull, 'w') as devnull:
                with redirect_stderr(devnull):
                    reader = PdfReader(path)
                    for i, page in enumerate(reader.pages):
                        try:
                            text = page.extract_text() or ""
                            text = text.strip()
                            if text:
                                texts.append({"text": text, "metadata": {"file": os.path.basename(path), "page": i + 1, "file_type": "pdf"}})
                        except Exception as page_error:
                            # Skip problematic pages but continue with others
                            print(f"Warning: Failed to extract text from page {i + 1} of {os.path.basename(path)} with pypdf: {page_error}")
                            continue
        except Exception as pypdf_error:
            print(f"Warning: pypdf also failed for {os.path.basename(path)}: {pypdf_error}")

    # If we got very little text overall, treat as scanned and OCR
    total_chars = sum(len(p["text"]) for p in texts)
    if total_chars < MIN_PDF_TEXT_LEN:
        return extract_from_scanned_pdf(path)
    return texts

def extract_from_scanned_pdf(path: str) -> List[Dict]:
    out = []
    try:
        pages = convert_from_path(path, dpi=200)  # requires poppler
        for i, page in enumerate(pages):
            try:
                text = pytesseract.image_to_string(page)
                text = text.strip()
                if text:
                    out.append({"text": text, "metadata": {"file": os.path.basename(path), "page": i + 1, "file_type": "scanned_pdf"}})
            except Exception as ocr_error:
                print(f"Warning: OCR failed for page {i + 1} of {os.path.basename(path)}: {ocr_error}")
                continue
    except Exception as convert_error:
        print(f"Warning: Failed to convert PDF to images for OCR: {os.path.basename(path)}: {convert_error}")
    return out

def extract_from_docx(path: str) -> List[Dict]:
    out = []
    try:
        doc = docx.Document(path)
        para_texts = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                para_texts.append(t)
        # Group paragraphs into larger pieces to avoid tiny chunks
        if para_texts:
            out.append({"text": "\n\n".join(para_texts), "metadata": {"file": os.path.basename(path), "file_type": "docx"}})
    except Exception:
        pass
    return out

def extract_from_pptx(path: str) -> List[Dict]:
    out = []
    try:
        prs = Presentation(path)
        for s_idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                out.append({"text": "\n".join(slide_text), "metadata": {"file": os.path.basename(path), "slide": s_idx + 1, "file_type": "pptx"}})
    except Exception:
        pass
    return out

def extract_from_html(path: str) -> List[Dict]:
    out = []
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            markup = f.read()
        soup = BeautifulSoup(markup, "html.parser")
        parts = []
        for el in soup.find_all(["h1", "h2", "h3", "p", "li"]):
            txt = el.get_text(separator=" ", strip=True)
            if txt:
                parts.append(txt)
        if parts:
            out.append({"text": "\n\n".join(parts), "metadata": {"file": os.path.basename(path), "file_type": "html"}})
    except Exception:
        pass
    return out

def extract_from_image(path: str) -> List[Dict]:
    out = []
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        text = text.strip()
        if text:
            out.append({"text": text, "metadata": {"file": os.path.basename(path), "file_type": "image"}})
    except Exception:
        pass
    return out

def extract_from_csv(path: str) -> List[Dict]:
    out = []
    try:
        df = pd.read_csv(path)
        rows = []
        for i, row in df.iterrows():
            parts = [f"{col}: {row[col]}" for col in df.columns]
            rows.append("; ".join(parts))
        if rows:
            out.append({"text": "\n".join(rows), "metadata": {"file": os.path.basename(path), "file_type": "csv"}})
    except Exception:
        pass
    return out


def extract_from_email(path: str) -> List[Dict]:
    """
    Extract text from email files (.eml, .emlx formats).
    
    Returns structured text with email metadata (from, to, subject, date) and body content.
    Handles both plain text and HTML emails, stripping HTML tags for indexing.
    """
    out = []
    try:
        # Read the email file
        with open(path, 'rb') as f:
            # Mac .emlx files have a header line with message length, skip it
            first_line = f.readline()
            # If it looks like a length header (just digits), it's .emlx format
            if first_line.strip().isdigit():
                # Continue reading from current position (after the length line)
                msg = BytesParser(policy=policy.default).parse(f)
            else:
                # Regular .eml file, rewind and parse from beginning
                f.seek(0)
                msg = BytesParser(policy=policy.default).parse(f)
        
        # Extract metadata
        subject = msg.get('subject', '(No Subject)')
        from_addr = msg.get('from', '(Unknown Sender)')
        to_addr = msg.get('to', '(Unknown Recipient)')
        date = msg.get('date', '(No Date)')
        
        # Extract body content
        body_text = ""
        if msg.is_multipart():
            # Handle multipart messages (plain text + HTML, attachments, etc.)
            for part in msg.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get('Content-Disposition', ''))
                
                # Skip attachments
                if 'attachment' in content_disposition:
                    continue
                    
                # Get plain text parts
                if content_type == 'text/plain':
                    try:
                        text = part.get_content()
                        if text:
                            body_text += text + "\n"
                    except Exception:
                        pass
                        
                # Get HTML parts and convert to text
                elif content_type == 'text/html' and not body_text:
                    # Only use HTML if we don't have plain text
                    try:
                        html_content = part.get_content()
                        # Use BeautifulSoup to extract text from HTML
                        soup = BeautifulSoup(html_content, 'html.parser')
                        # Remove script and style elements
                        for script in soup(['script', 'style']):
                            script.decompose()
                        text = soup.get_text(separator='\n', strip=True)
                        if text:
                            body_text += text + "\n"
                    except Exception:
                        pass
        else:
            # Simple non-multipart message
            content_type = msg.get_content_type()
            if content_type == 'text/plain':
                try:
                    body_text = msg.get_content()
                except Exception:
                    pass
            elif content_type == 'text/html':
                try:
                    html_content = msg.get_content()
                    soup = BeautifulSoup(html_content, 'html.parser')
                    for script in soup(['script', 'style']):
                        script.decompose()
                    body_text = soup.get_text(separator='\n', strip=True)
                except Exception:
                    pass
        
        # Clean up the body text
        body_text = body_text.strip()
        
        # Construct the final text for indexing with metadata
        email_text_parts = [
            f"Subject: {subject}",
            f"From: {from_addr}",
            f"To: {to_addr}",
            f"Date: {date}",
            "",
            body_text
        ]
        
        final_text = "\n".join(email_text_parts)
        
        if final_text.strip():
            out.append({
                "text": final_text,
                "metadata": {
                    "file": os.path.basename(path),
                    "file_type": "email",
                    "subject": subject,
                    "from": from_addr,
                    "to": to_addr,
                    "date": date
                }
            })
    except Exception as e:
        print(f"Error extracting email from {path}: {e}")
        pass
    
    return out


def extract_from_audio(path: str) -> List[Dict]:
    """
    Extract text from audio files via Whisper API transcription.
    
    Supports: .wav, .mp3, .m4a, .flac, .ogg
    Returns sentence-chunked text for better semantic coherence.
    """
    from app.config import settings
    
    out = []
    try:
        filename = os.path.basename(path)
        file_size = os.path.getsize(path)
        # logger.info(f"[AUDIO] Starting transcription: {filename} (size: {file_size} bytes)")
        # logger.info(f"[AUDIO] Whisper API URL: {settings.whisper_api_url}")
        # logger.info(f"[AUDIO] Language: {settings.whisper_language}, Timeout: {settings.whisper_api_timeout}s")
        
        # Prepare the file for upload
        with open(path, 'rb') as audio_file:
            files = {'file': (filename, audio_file, 'audio/wav')}
            data = {
                'language': settings.whisper_language,
                'task': 'transcribe'
            }
            
            # logger.info(f"[AUDIO] Sending request to Whisper API...")
            # Call Whisper API
            response = requests.post(
                settings.whisper_api_url,
                files=files,
                data=data,
                timeout=settings.whisper_api_timeout
            )
            # logger.info(f"[AUDIO] Received response: status={response.status_code}")
            response.raise_for_status()
        
        # Parse JSON response: {"text": str, "language": str, "duration": float}
        # logger.info(f"[AUDIO] Parsing JSON response...")
        result = response.json()
        transcribed_text = result.get('text', '').strip()
        detected_language = result.get('language', settings.whisper_language)
        duration = result.get('duration', 0)
        
        # logger.info(f"[AUDIO] Transcription received: {len(transcribed_text)} chars, language={detected_language}, duration={duration}s")
        
        if not transcribed_text:
            logger.warning(f"[AUDIO] Warning: No text transcribed from {filename}")
            return []
        
        # Split into sentences for better chunking
        # logger.info(f"[AUDIO] Tokenizing into sentences...")
        sentences = sent_tokenize(transcribed_text)
        # logger.info(f"[AUDIO] Found {len(sentences)} sentences")
        
        # Group sentences into reasonable chunks (avoid single-sentence chunks)
        # Aim for ~3-5 sentences per chunk
        chunk_size = 4
        for i in range(0, len(sentences), chunk_size):
            chunk_sentences = sentences[i:i + chunk_size]
            chunk_text = ' '.join(chunk_sentences)
            
            out.append({
                "text": chunk_text,
                "metadata": {
                    "file": filename,
                    "file_type": "audio",
                    "language": detected_language,
                    "duration": duration,
                    "chunk_id": i // chunk_size
                }
            })
        
        # logger.info(f"[AUDIO] Created {len(out)} chunks from {filename}")
        # logger.info(f"[AUDIO] SUCCESS: Transcribed {filename}: {len(sentences)} sentences in {len(out)} chunks")
        
    except requests.exceptions.Timeout:
        logger.error(f"[AUDIO] ERROR: Whisper API timeout for {os.path.basename(path)} (>{settings.whisper_api_timeout}s)")
    except requests.exceptions.RequestException as e:
        logger.error(f"[AUDIO] ERROR: Whisper API request failed for {os.path.basename(path)}: {e}")
    except Exception as e:
        import traceback
        logger.error(f"[AUDIO] ERROR: Exception transcribing audio {os.path.basename(path)}: {e}")
        logger.error(f"[AUDIO] Traceback: {traceback.format_exc()}")
    
    return out


def extract_fallback_text(path: str) -> List[Dict]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read().strip()
        if text:
            return [{"text": text, "metadata": {"file": os.path.basename(path), "file_type": "text"}}]
    except Exception:
        pass
    return []

def extract(path: str) -> List[Dict]:
    mime = detect_mime(path) or ""
    lower = path.lower()
    filename = os.path.basename(path)
    print(f"Detected MIME type for {filename}: {mime}")
    try:
        if lower.endswith(".pdf") or (mime and mime.startswith("application/pdf")):
            print(f"Processing PDF: {filename}")
            return extract_from_pdf(path)
        if lower.endswith(".docx"):
            print(f"Processing DOCX: {filename}")
            return extract_from_docx(path)
        if lower.endswith(".pptx"):
            print(f"Processing PPTX: {filename}")
            return extract_from_pptx(path)
        if lower.endswith(".html") or lower.endswith(".htm") or mime == "text/html":
            print(f"Processing HTML: {filename}")
            return extract_from_html(path)
        if lower.endswith(".csv"):
            print(f"Processing CSV: {filename}")
            return extract_from_csv(path)
        if lower.endswith((".eml", ".emlx")):
            print(f"Processing Email: {filename}")
            return extract_from_email(path)
        if lower.endswith((".wav", ".mp3", ".m4a", ".flac", ".ogg")) or (mime and mime.startswith("audio/")):
            print(f"Processing Audio: {filename}")
            return extract_from_audio(path)
        if lower.endswith((".png", ".jpg", ".jpeg", ".tiff", ".tif")) or (mime and mime.startswith("image/")):
            print(f"Processing Image: {filename}")
            return extract_from_image(path)
        # fallback: try plain text
        print(f"Processing as text: {filename}")
        return extract_fallback_text(path)
    except Exception as e:
        print(f"Error processing file {filename}: {e}")
        return []